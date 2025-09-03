# core/services.py

import os
import uuid
from collections import OrderedDict
from dotenv import load_dotenv
import google.generativeai as genai

# For file processing
import PyPDF2
import docx
import openpyxl
from pptx import Presentation

# For AI and vector store
from sentence_transformers import SentenceTransformer
import faiss
import numpy as np

# --- NEW: Use a dictionary to hold data for MULTIPLE user sessions ---
# The key will be the user's session ID.
# The value will be their personal OrderedDict of documents.
multi_session_data = {}
MAX_DOCUMENTS = 3

# --- Load environment variables and AI models once ---
load_dotenv()
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))
embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
llm = genai.GenerativeModel('gemini-1.5-flash')


def get_user_session(session_id):
    """Gets or creates the document store for a specific user session."""
    if session_id not in multi_session_data:
        multi_session_data[session_id] = OrderedDict()
    return multi_session_data[session_id]

def get_recent_documents(session_id):
    """Returns a list of recent documents for a specific user."""
    user_session = get_user_session(session_id)
    return [{"id": doc_id, "name": data["name"]} for doc_id, data in user_session.items()]

def process_document(file, session_id):
    """Processes a document for a specific user session."""
    user_session = get_user_session(session_id)
    
    try:
        document_id = str(uuid.uuid4())
        file_name = file.name
        file.seek(0)
        text = ""

        # Logic to handle different file types
        file_extension = os.path.splitext(file_name)[1].lower()

        if file_extension == '.pdf':
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() or ""
        
        elif file_extension == '.docx':
            doc = docx.Document(file)
            for para in doc.paragraphs:
                text += para.text + "\n"

        elif file_extension == '.xlsx':
            workbook = openpyxl.load_workbook(file)
            for sheet in workbook.worksheets:
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value:
                            text += str(cell.value) + " "
                    text += "\n"
        
        elif file_extension == '.pptx':
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        
        else:
            print(f"Unsupported file type: {file_extension}")
            return None, None
        
        chunks = [p.strip() for p in text.split("\n\n") if p.strip()]
        if not chunks:
            print("No text could be extracted from the document.")
            return None, None
            
        embeddings = embedding_model.encode(chunks)
        
        index = faiss.IndexFlatL2(embeddings.shape[1])
        index.add(embeddings.astype('float32'))
        
        if len(user_session) >= MAX_DOCUMENTS:
            user_session.popitem(last=False)

        user_session[document_id] = {
            "chunks": chunks,
            "index": index,
            "name": file_name
        }
        
        return document_id, file_name
    except Exception as e:
        print(f"Error processing document: {e}")
        return None, None

def query_document(document_id, question, session_id):
    """Queries a document from a specific user session."""
    user_session = get_user_session(session_id)
    if document_id not in user_session:
        return "Document not found. It may have been cleared from the session. Please upload it again."

    document_data = user_session[document_id]
    index = document_data["index"]
    chunks = document_data["chunks"]
    
    question_embedding = embedding_model.encode([question])
    
    D, I = index.search(question_embedding.astype('float32'), k=3)
    
    relevant_chunks = [chunks[i] for i in I[0]]
    context = "\n\n".join(relevant_chunks)
    
    prompt = f"""
    You are a helpful Q&A assistant. Use the following pieces of context from a document to answer the user's question.
    If the answer is not contained within the provided text, state that you cannot find the answer in the document.

    Context:
    {context}

    Question: {question}

    Answer:
    """
    
    try:
        response = llm.generate_content(prompt)
        return response.text
    except Exception as e:
        print(f"Error generating content: {e}")
        return f"An error occurred while generating the answer: {e}"